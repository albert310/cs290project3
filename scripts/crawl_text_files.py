#!/usr/bin/env python3
"""
Scan all crawled markdown files for text-type file links, download them,
and convert to markdown where possible (.docx, .doc, .xlsx, .xls, .txt, .csv, etc.)
"""
import asyncio
import re
import sys
import time
import json
import hashlib
import subprocess
import urllib.parse
from pathlib import Path
from urllib.parse import urlparse, urljoin

import aiohttp

CRAWL_DIR = Path("/Users/leslie/Desktop/data/shanghaitech")
OUTPUT_DIR = Path("/Users/leslie/Desktop/data/shanghaitech_text")
OUTPUT_DIR_MD = Path("/Users/leslie/Desktop/data/shanghaitech_text_md")
STATE_FILE = Path("/Users/leslie/Desktop/data/crawl_text_state.json")

# File extensions to target
TARGET_EXTS = {
    "txt", "csv", "xml", "json", "tex", "rtf", "log", "md", "rst", "yaml", "yml",
    "doc", "docx", "xls", "xlsx", "ppt", "pptx", "odt", "ods", "odp",
}

# Regex to find file URLs in markdown
LINK_RE = re.compile(
    r'(?:href|src|url)\s*=\s*["\']([^"\']+\.(' + "|".join(TARGET_EXTS) + r'))["\']|'
    r'\[([^\]]*)\]\(([^)]+\.(' + "|".join(TARGET_EXTS) + r'))\)|'
    r'(?<!\()(https?://[^\s<>"\')\]]+\.(' + "|".join(TARGET_EXTS) + r'))(?![\w-])',
    re.IGNORECASE
)


def extract_file_urls(text, base_url):
    """Extract all text-type file URLs from markdown text."""
    found = set()

    # Match markdown links: [text](url)
    for m in re.finditer(r'\[.*?\]\((https?://[^\)]+)\)', text):
        url = m.group(1)
        parsed = urlparse(url)
        # Strip query strings to check extension
        path_lower = parsed.path.lower()
        for ext in TARGET_EXTS:
            if path_lower.endswith("." + ext):
                found.add(url)
                break

    # Match plain URLs
    for m in re.finditer(r'(?<!\()(https?://[^\s<>"\')\]]+)', text):
        url = m.group(0).rstrip(".,;:!?\"'")
        parsed = urlparse(url)
        path_lower = parsed.path.lower()
        for ext in TARGET_EXTS:
            if path_lower.endswith("." + ext):
                found.add(url)
                break

    # Match href/src attributes
    for m in re.finditer(r'(?:href|src)\s*=\s*["\']([^"\']+)["\']', text):
        path = m.group(1)
        parsed = urlparse(path)
        path_lower = parsed.path.lower()
        for ext in TARGET_EXTS:
            if path_lower.endswith("." + ext):
                if path.startswith("http"):
                    found.add(path)
                else:
                    try:
                        found.add(urljoin(base_url, path))
                    except Exception:
                        pass
                break

    return found


def url_to_relpath(url):
    """Convert URL to a relative filesystem path."""
    p = urlparse(url)
    netloc = p.netloc.replace(":", "_")
    path = p.path.strip("/") or "index"

    # Keep directories structure
    parts = path.split("/")
    safe_parts = []
    for part in parts:
        part = re.sub(r'[<>:"/\\|?*\s]+', "_", part)
        if len(part) > 100:
            part = part[:100]
        safe_parts.append(part)
    safe_path = "/".join(safe_parts)

    if len(safe_path) > 200:
        safe_path = safe_path[:200]

    # Use query hash to avoid overwriting same-name files
    if p.query:
        qhash = hashlib.md5(p.query.encode()).hexdigest()[:8]
        base, ext = safe_path.rsplit(".", 1) if "." in safe_path.split("/")[-1] else (safe_path, "")
        if ext:
            safe_path = f"{base}__{qhash}.{ext}"

    return f"{netloc}/{safe_path}"


def load_state():
    if STATE_FILE.exists():
        with open(STATE_FILE, "r") as f:
            return json.load(f)
    return {"found_urls": [], "downloaded": {}, "failed": {}}


def save_state(state):
    with open(STATE_FILE, "w") as f:
        json.dump(state, f, indent=2, ensure_ascii=False)


def scan_all_markdown():
    """Scan all markdown files for text-type file URLs."""
    print("Scanning crawled markdown for text-type file links...")
    all_urls = {}
    file_count = 0

    for md_file in CRAWL_DIR.rglob("*.md"):
        file_count += 1
        try:
            text = md_file.read_text(encoding="utf-8")
            # Build base URL from file path
            rel = md_file.relative_to(CRAWL_DIR)
            domain = str(rel).split("/")[0]
            base_url = f"https://{domain}"
            urls = extract_file_urls(text, base_url)
            for url in urls:
                if url not in all_urls:
                    all_urls[url] = str(md_file)
        except Exception:
            pass

    print(f"  Scanned {file_count} files, found {len(all_urls)} unique text-type file URLs")
    return all_urls


async def download_file(session, url, filepath, sem):
    """Download a single file."""
    async with sem:
        try:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=60)) as resp:
                if resp.status == 200:
                    content = await resp.read()
                    filepath.parent.mkdir(parents=True, exist_ok=True)
                    filepath.write_bytes(content)
                    return True, len(content), None
                else:
                    return False, 0, f"HTTP {resp.status}"
        except asyncio.TimeoutError:
            return False, 0, "timeout"
        except Exception as e:
            return False, 0, str(e)[:100]


def convert_to_markdown(src_path, dst_path):
    """Convert downloaded file to markdown.

    Handles:
    - .txt/.csv/.xml/.json/.tex/.rtf/.md/.rst/.yaml/.yml/.log -> direct copy as .md
    - .docx -> pandoc to markdown
    - .doc -> libreoffice/pandoc
    - .xlsx/.xls -> pandas/ssconvert to CSV then save as markdown table
    - .pptx -> python-pptx extract text
    - .odt -> pandoc
    """
    suffix = src_path.suffix.lower()
    dst_path.parent.mkdir(parents=True, exist_ok=True)

    # Plain text formats: just copy with .md extension
    plain_exts = {".txt", ".csv", ".xml", ".json", ".tex", ".rtf", ".md", ".rst", ".yaml", ".yml", ".log"}
    if suffix in plain_exts:
        try:
            text = src_path.read_text(encoding="utf-8", errors="replace")
            dst_path.write_text(text, encoding="utf-8")
            return True
        except Exception as e:
            print(f"    read error: {e}")
            return False

    # docx -> pandoc
    if suffix == ".docx":
        try:
            subprocess.run(
                ["pandoc", str(src_path), "-f", "docx", "-t", "markdown", "-o", str(dst_path)],
                capture_output=True, timeout=30
            )
            if dst_path.exists():
                return True
        except Exception:
            pass
        # Fallback: python-docx
        try:
            from docx import Document
            doc = Document(str(src_path))
            paragraphs = [p.text for p in doc.paragraphs]
            dst_path.write_text("\n\n".join(paragraphs), encoding="utf-8")
            return True
        except ImportError:
            pass
        except Exception as e:
            print(f"    docx error: {e}")
        return False

    # doc -> pandoc or textract
    if suffix == ".doc":
        try:
            subprocess.run(
                ["pandoc", str(src_path), "-f", "doc", "-t", "markdown", "-o", str(dst_path)],
                capture_output=True, timeout=30
            )
            if dst_path.exists():
                return True
        except Exception:
            pass
        # Fallback: antiword or catdoc
        for tool in ["antiword", "catdoc", "textutil"]:
            try:
                if tool == "textutil":
                    result = subprocess.run(
                        ["textutil", "-convert", "txt", "-stdout", str(src_path)],
                        capture_output=True, text=True, timeout=30
                    )
                else:
                    result = subprocess.run(
                        [tool, str(src_path)], capture_output=True, text=True, timeout=30
                    )
                if result.stdout.strip():
                    dst_path.write_text(result.stdout, encoding="utf-8")
                    return True
            except Exception:
                continue
        return False

    # xlsx/xls -> ssconvert to CSV, then markdown table
    if suffix in {".xlsx", ".xls"}:
        # Try pandas first
        try:
            import pandas as pd
            df = pd.read_excel(str(src_path), engine="openpyxl" if suffix == ".xlsx" else "xlrd")
            # Convert to markdown table
            md = df.to_markdown(index=False) if hasattr(df, "to_markdown") else df.to_string(index=False)
            dst_path.write_text(md, encoding="utf-8")
            return True
        except ImportError:
            pass
        except Exception as e:
            print(f"    pandas error: {e}")

        # Fallback: python openpyxl/xlrd
        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(src_path), data_only=True)
                output = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    output.append(f"## Sheet: {sheet_name}\n")
                    for row in ws.iter_rows(values_only=True):
                        output.append(" | ".join(str(c) if c is not None else "" for c in row))
                    output.append("")
                dst_path.write_text("\n".join(output), encoding="utf-8")
                return True
            except ImportError:
                pass
            except Exception:
                pass
        return False

    # pptx -> extract text
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(src_path))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"## Slide {slide_num}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                text_parts.append("")
            dst_path.write_text("\n\n".join(text_parts), encoding="utf-8")
            return True
        except ImportError:
            pass
        except Exception as e:
            print(f"    pptx error: {e}")
        return False

    # odt -> pandoc
    if suffix == ".odt":
        try:
            subprocess.run(
                ["pandoc", str(src_path), "-f", "odt", "-t", "markdown", "-o", str(dst_path)],
                capture_output=True, timeout=30
            )
            if dst_path.exists():
                return True
        except Exception:
            pass
        return False

    return False


async def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR_MD.mkdir(parents=True, exist_ok=True)

    # Step 1: Scan all markdown for text file URLs
    url_map = scan_all_markdown()
    state = load_state()

    # Filter already downloaded
    todo_urls = {url: src for url, src in url_map.items() if url not in state["downloaded"]}
    print(f"  Already downloaded: {len(state['downloaded'])}")
    print(f"  Still need to download: {len(todo_urls)}")

    if not todo_urls:
        print("All text files already downloaded!")
        return

    # Step 2: Download with concurrency
    print(f"\nDownloading {len(todo_urls)} files...")
    sem = asyncio.Semaphore(10)
    success_count = 0
    fail_count = 0

    connector = aiohttp.TCPConnector(limit=20, force_close=True)
    async with aiohttp.ClientSession(connector=connector) as session:
        tasks = []
        url_list = list(todo_urls.keys())
        for url in url_list:
            relpath = url_to_relpath(url)
            filepath = OUTPUT_DIR / relpath
            tasks.append(download_file(session, url, filepath, sem))

        results = await asyncio.gather(*tasks)

    for url, (ok, size, err) in zip(url_list, results):
        if ok:
            state["downloaded"][url] = {"size": size, "time": time.time()}
            success_count += 1
        else:
            state["failed"][url] = {"error": err, "time": time.time()}
            fail_count += 1

    save_state(state)
    print(f"  Downloaded: {success_count}, Failed: {fail_count}")

    # Step 3: Convert to markdown
    print(f"\nConverting downloaded files to markdown...")
    converted = 0
    skipped = 0

    for url in state["downloaded"]:
        relpath = url_to_relpath(url)
        src_path = OUTPUT_DIR / relpath
        # Same path structure but under _md dir
        md_relpath = str(relpath)
        # Change extension to .md
        base = md_relpath.rsplit(".", 1)[0] if "." in md_relpath.split("/")[-1] else md_relpath
        dst_path = OUTPUT_DIR_MD / (base + ".md")

        if dst_path.exists():
            skipped += 1
            continue

        ok = convert_to_markdown(src_path, dst_path)
        if ok:
            converted += 1

    print(f"  Converted: {converted}, Skipped (already exist): {skipped}")

    # Summary
    print(f"\n=== Summary ===")
    print(f"Total text-type file URLs found: {len(url_map)}")
    print(f"Files downloaded: {len(state['downloaded'])}")
    print(f"Files failed: {len(state['failed'])}")
    print(f"Files at: {OUTPUT_DIR}")
    print(f"Markdown at: {OUTPUT_DIR_MD}")

    # Extension breakdown
    ext_counts = {}
    for url in url_map:
        parsed = urlparse(url)
        ext = parsed.path.rsplit(".", 1)[-1].lower() if "." in parsed.path else "unknown"
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    print(f"\nFile type breakdown:")
    for ext, cnt in sorted(ext_counts.items(), key=lambda x: -x[1]):
        print(f"  .{ext}: {cnt}")


if __name__ == "__main__":
    asyncio.run(main())
